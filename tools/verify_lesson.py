"""Check a finished lesson before it reaches a student.

    python tools/verify_lesson.py 05          # the mechanical checks
    python tools/verify_lesson.py 05 --run    # also execute every notebook
    python tools/verify_lesson.py             # every lesson that exists

The rule this exists to enforce: **a student should never be the one who finds
the error.**

What it can and cannot do is worth being precise about, because the failure that
prompted it would have slipped through a naive version. Lesson 3 once claimed
that X-transpose y was 165,200 when the data gives 165,600. Every number after
it — the determinant, the inverse, the two components of theta — was derived
*correctly from the wrong figure*, so the example was perfectly self-consistent
and read as careful work. No consistency check catches that. Only recomputing
from the raw data does.

So there are two kinds of check here, and only one of them is automatic:

**Mechanical** — the checks in this file. Notebooks that run, figures that
exist, formulas that survive conversion, slide counts that match the lesson
plan, acronyms expanded, cross-references that resolve, carry-home numbers
quoted from another lesson that still match it. These run on every lesson,
every time, and cost nothing.

**Arithmetic** — a per-lesson ``Docs/worked_examples.py``. Every number a
handout works out by hand must be recomputed there from the raw inputs, by a
route that does not reuse the handout's intermediate values, and asserted
against what the handout prints. This file discovers and runs those; it cannot
write them, and a lesson that works numbers by hand without one is reported as
incomplete.

Exit code is 0 when everything passes, 1 otherwise, so it can gate a commit.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from functools import lru_cache
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
CONTAINER = "tai_course"

DOLLAR = chr(36)
DISPLAY_MATH = re.compile(re.escape(DOLLAR * 2) + "(.+?)" + re.escape(DOLLAR * 2), re.S)
INLINE_MATH = re.compile(
    "(?<!" + re.escape(DOLLAR) + ")" + re.escape(DOLLAR)
    + "([^" + DOLLAR + "\n]+?)" + re.escape(DOLLAR)
    + "(?!" + re.escape(DOLLAR) + ")")
FIGURE_REF = re.compile(r"!\[[^\]]*\]\(([^)]+)\)")

#: Acronyms the course expands on first use in every artefact. Kept here rather
#: than in a doc so the check and the rule cannot drift apart.
ACRONYMS = {
    "MSE": ("mean squared error",),
    "RMSE": ("root mean squared error",),
    "MCAR": ("missing completely at random",),
    "MAR": ("missing at random",),
    "MNAR": ("missing not at random", "not missing at random"),
    "IQR": ("interquartile range",),
    "OLS": ("ordinary least squares",),
    "AUC": ("area under",),
    "ROC": ("receiver operating",),
    "PCA": ("principal component",),
    "SVM": ("support vector",),
    "EDA": ("exploratory data analysis",),
    "API": ("application programming interface",),
    "MLE": ("maximum likelihood",),
    "SMART": ("self monitoring",),
    "TPR": ("true positive rate",),
    "FPR": ("false positive rate",),
}


class Report:
    """Collects problems so one run reports all of them, not just the first."""

    def __init__(self) -> None:
        self.problems: list[str] = []
        self.notes: list[str] = []

    def fail(self, check: str, detail: str) -> None:
        self.problems.append(f"{check}: {detail}")

    def note(self, detail: str) -> None:
        self.notes.append(detail)

    @property
    def ok(self) -> bool:
        return not self.problems


def markdown_of(path: Path) -> str:
    """The prose of a file: markdown cells only, for a notebook."""
    if path.suffix == ".ipynb":
        nb = json.loads(path.read_text(encoding="utf8"))
        return "\n".join("".join(c.get("source", []))
                         for c in nb.get("cells", [])
                         if c.get("cell_type") == "markdown")
    return path.read_text(encoding="utf8")


def shipped(paths) -> list[Path]:
    """Those of `paths` that a student will actually receive.

    Instructor-only content is marked by a git ignore rule: the worked
    exercise solutions, generated for whoever marks the work and never
    committed. They are held to the same standards while being written, but
    not by this gate, which exists to check what reaches a student. A red
    result on a file no student will open is a red result people learn to
    scroll past.

    Git is asked about the paths in hand, on every call, rather than once into
    a cache at start-up. A `--run` pass spends twenty minutes executing
    notebooks, which is long enough for an instructor file to be written while
    it is still going; a set built before that file existed reports it as
    shippable. That is how lesson 9's solution notebook came to be checked,
    and failed, for an acronym no student will ever read.
    """
    paths = list(paths)
    if not paths:
        return []
    try:
        # check-ignore exits 1 when nothing matches, which is not an error,
        # so this deliberately does not pass check=True.
        result = subprocess.run(
            ["git", "check-ignore", "--stdin"], cwd=ROOT,
            capture_output=True, text=True,
            input="\n".join(str(p) for p in paths))
    except (FileNotFoundError, OSError):
        return paths                      # not a checkout: check everything
    skip = {Path(line).resolve() for line in result.stdout.splitlines() if line}
    return [p for p in paths if p.resolve() not in skip]


def artefacts(lesson: Path) -> list[Path]:
    return shipped(p for p in sorted(lesson.rglob("*"))
                   if p.suffix in (".md", ".ipynb")
                   and "checkpoint" not in str(p)
                   and "__pycache__" not in str(p))


# --------------------------------------------------------------- notebooks

def check_notebooks(lesson: Path, report: Report, run: bool) -> None:
    notebooks = sorted((lesson / "Notebooks").glob("*.ipynb"))
    notebooks += sorted((lesson / "Quizzes").glob("*.ipynb"))
    notebooks += shipped(sorted((lesson / "Exercises").glob("*.ipynb")))
    if not notebooks:
        report.fail("notebooks", "none found")
        return

    for path in notebooks:
        name = path.name
        try:
            nb = json.loads(path.read_text(encoding="utf8"))
        except json.JSONDecodeError as error:
            report.fail("notebooks", f"{name} is not valid JSON: {error}")
            continue

        code_cells = [c for c in nb["cells"] if c["cell_type"] == "code"]
        is_quiz = "Quizzes" in path.parts

        if is_quiz:
            if code_cells:
                report.fail("quiz", f"{name} has {len(code_cells)} code cells; "
                                    "quizzes are markdown only")
            continue

        if not code_cells:
            report.fail("notebooks", f"{name} has no code cells")
            continue

        # A stored traceback means the committed notebook is broken.
        for index, cell in enumerate(code_cells, 1):
            for output in cell.get("outputs", []):
                if output.get("output_type") == "error":
                    report.fail("notebooks",
                                f"{name} cell {index} stored an error: "
                                f"{output.get('ename')}")

        # Outputs are committed on purpose (see CLAUDE.md); missing ones mean
        # the notebook was edited and never re-run.
        without = [i for i, c in enumerate(code_cells, 1) if not c.get("outputs")]
        if len(without) > 1:
            report.fail("notebooks",
                        f"{name}: {len(without)} code cells have no output "
                        f"(cells {without[:6]}) — re-run before committing")

        # The pinned versions must be the ones the image actually ships.
        source = "".join("".join(c["source"]) for c in code_cells)
        if "pip install" not in source:
            report.fail("notebooks", f"{name} has no pinned pip install line")

    if run:
        for path in notebooks:
            if "Quizzes" in path.parts:
                continue
            relative = path.relative_to(ROOT)
            # Execute into a scratch directory, never over the committed file:
            # --inplace would rewrite execution counts and kernel metadata, so
            # checking a lesson would leave every notebook modified.
            result = subprocess.run(
                ["docker", "exec", "-w", f"/home/jovyan/work/{relative.parent}",
                 CONTAINER, "jupyter", "nbconvert", "--to", "notebook",
                 "--execute", "--output-dir", "/tmp/verify",
                 "--ExecutePreprocessor.timeout=1800", path.name],
                capture_output=True, text=True)
            if result.returncode != 0:
                tail = result.stderr.strip().splitlines()[-3:]
                report.fail("notebooks --run",
                            f"{path.name} failed: {' | '.join(tail)}")
            else:
                report.note(f"runs clean: {path.name}")


def check_pins(lesson: Path, report: Report) -> None:
    """Pinned versions must match what the container actually has."""
    try:
        result = subprocess.run(
            ["docker", "exec", CONTAINER, "python", "-c",
             "import numpy, pandas, sklearn, matplotlib; "
             "print(numpy.__version__, pandas.__version__, "
             "sklearn.__version__, matplotlib.__version__)"],
            capture_output=True, text=True)
    except FileNotFoundError:
        report.note("docker not installed; skipped the pinned-version check")
        return
    if result.returncode != 0:
        report.note("container not running; skipped the pinned-version check")
        return

    numpy_v, pandas_v, sklearn_v, matplotlib_v = result.stdout.split()
    expected = {"numpy": numpy_v, "pandas": pandas_v,
                "scikit-learn": sklearn_v, "matplotlib": matplotlib_v}

    for path in artefacts(lesson):
        if path.suffix != ".ipynb":
            continue
        text = path.read_text(encoding="utf8")
        for package, version in expected.items():
            for found in re.findall(package + r"==([0-9][0-9.]*)", text):
                if found != version:
                    report.fail("pinned versions",
                                f"{path.name} pins {package}=={found}, "
                                f"the image has {version}")


# ----------------------------------------------------------------- figures

def check_figures_are_tracked(figures_dir: Path, report: Report) -> None:
    """A figure on this disk and not in the index is a figure nobody else has.

    Everything else here asks the filesystem, which is the one place a figure
    is certain to exist: it is sitting there because this machine just built
    it. The deck embeds it, the PDF renders it, every check passes, and the
    repository ships a reference to a file it does not carry.

    Twice in a day on 30 August, both times an equation image whose name
    changed with its LaTeX - `git add` on the modified paths does not stage a
    new one. Neither was caught by anything but reading `git status` by eye.
    """
    result = subprocess.run(
        ["git", "ls-files", "--", str(figures_dir)],
        capture_output=True, text=True, cwd=str(ROOT))
    if result.returncode != 0:
        return  # not a git checkout; nothing to say
    tracked = {Path(line).name for line in result.stdout.splitlines()}
    for path in sorted(figures_dir.glob("*.png")):
        if path.name not in tracked:
            report.fail("figures",
                        f"{path.name} is not tracked by git - it exists here "
                        "and nowhere else; git add it")


def check_figures(lesson: Path, report: Report) -> None:
    figures_dir = lesson / "Figures"
    if not figures_dir.exists():
        return

    referenced: set[str] = set()
    for path in artefacts(lesson):
        for name in FIGURE_REF.findall(markdown_of(path)):
            referenced.add(Path(name).name)
            if not (figures_dir / Path(name).name).exists():
                report.fail("figures",
                            f"{path.name} references {name}, which is missing")

    on_disk = {p.name for p in figures_dir.glob("*.png")}
    check_figures_are_tracked(figures_dir, report)

    # Equation images are generated; an unused one is stale output.
    # `~$name.pptx` is PowerPoint's owner file, written while the deck is open.
    # It is 165 bytes and not a zip, and globbing it crashed this check.
    decks = (p for p in (lesson / "Slides").glob("*.pptx")
             if not p.name.startswith("~$"))
    deck = next(iter(decks), None)
    if deck is not None:
        embedded = set()
        with zipfile.ZipFile(deck) as archive:
            for entry in archive.namelist():
                if entry.startswith("ppt/media/"):
                    embedded.add(hashlib.sha256(archive.read(entry)).hexdigest())
        for path in sorted(figures_dir.glob("eq_*.png")):
            digest = hashlib.sha256(path.read_bytes()).hexdigest()
            if digest not in embedded:
                report.fail("figures", f"{path.name} is an orphaned equation image")

    # Data figures come from notebooks, conceptual ones from make_figures.py.
    produced = ""
    for path in sorted((lesson / "Notebooks").glob("*.ipynb")):
        produced += path.read_text(encoding="utf8")
    maker = figures_dir / "make_figures.py"
    if maker.exists():
        produced += maker.read_text(encoding="utf8")

    for name in sorted(on_disk):
        if name.startswith("eq_"):
            continue
        if name not in produced:
            report.fail("figures",
                        f"{name} is not produced by any notebook or "
                        "make_figures.py — its origin is unreproducible")
        elif name not in referenced:
            report.note(f"{name} is generated but never shown")


    # The handout is what gets read after the lecture, when the projected
    # figures are gone. Every figure the slides show, it shows too.
    slide_figures: set[str] = set()
    for path in sorted((lesson / "Slides").glob("*_slides.md")):
        slide_figures |= {Path(n).name for n in FIGURE_REF.findall(
            path.read_text(encoding="utf8"))}

    # Reading material means the handout or the supplementary Resources:
    # the history figures belong with the history essay, not crammed into
    # the one section of the handout that summarises it.
    handout_figures: set[str] = set()
    reading = sorted((lesson / "Docs").glob("*.md")) + sorted(
        (lesson / "Resources").glob("*.md"))
    for path in reading:
        handout_figures |= {Path(n).name for n in FIGURE_REF.findall(
            path.read_text(encoding="utf8"))}

    # Equation images are excluded: the handout renders real LaTeX.
    missing = {n for n in slide_figures - handout_figures
               if not n.startswith("eq_")}
    for name in sorted(missing):
        report.fail("handout figures",
                    f"{name} is on a slide but not in the handout")


    # Referencing a figure is not the same as shipping one. pandoc drops an
    # image it cannot find without a word of complaint, so look in the PDF.
    for handout in sorted((lesson / "Docs").glob("*.md")):
        referenced = {Path(n).name for n in FIGURE_REF.findall(
            handout.read_text(encoding="utf8"))}
        pdf = handout.with_suffix(".pdf")
        if not referenced or not pdf.exists():
            continue
        embedded = len(re.findall(rb"/Subtype\s*/Image", pdf.read_bytes()))
        if embedded < len(referenced):
            report.fail("handout figures",
                        f"{handout.name} references {len(referenced)} figures "
                        f"but {pdf.name} embeds {embedded} images — pandoc "
                        "could not find them")


# ------------------------------------------------------------------ slides

def check_slides(lesson: Path, report: Report) -> None:
    sources = sorted((lesson / "Slides").glob("*_slides.md"))
    if not sources:
        report.fail("slides", "no slides source found")
        return
    source = sources[0]
    text = source.read_text(encoding="utf8")

    # Inline maths must survive the Unicode conversion; there is no image
    # fallback inline, so anything else silently mangles on the slide.
    sys.path.insert(0, str(ROOT / "tools"))
    import render_math

    for match in INLINE_MATH.finditer(text):
        formula = match.group(1)
        rendered = render_math.to_unicode(formula)
        if rendered is None or chr(92) in str(rendered):
            report.fail("slides maths",
                        f"inline formula has no Unicode form: {DOLLAR}{formula}{DOLLAR}")

    deck = source.with_suffix(".pptx")
    if not deck.exists():
        report.fail("slides", f"{deck.name} has not been built")
        return

    from pptx import Presentation
    slides = list(Presentation(deck).slides)

    notes_by_slide: dict[int, str] = {}

    for index, slide in enumerate(slides, 1):
        has_text = any(sh.has_text_frame and sh.text_frame.text.strip()
                       for sh in slide.shapes)
        has_picture = any(sh.shape_type == 13 for sh in slide.shapes)
        has_table = any(sh.has_table for sh in slide.shapes)
        if not (has_text or has_picture or has_table):
            report.fail("slides", f"slide {index} is empty")
        if not slide.shapes.title or not slide.shapes.title.text.strip():
            report.fail("slides", f"slide {index} has no title")

        if index == 1:
            continue
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        notes_by_slide[index] = notes
        title = slide.shapes.title.text.strip() if slide.shapes.title else ""
        if len(notes) < 60 and title.lower() not in ("break",):
            report.fail("slides",
                        f"slide {index} ({title[:34]}) has thin speaker notes")

    check_lesson_plan(lesson, len(slides), notes_by_slide, report)


def check_lesson_plan(lesson: Path, slide_count: int,
                      notes_by_slide: dict, report: Report) -> None:
    handouts = sorted((lesson / "Docs").glob("*.md"))
    if not handouts:
        report.fail("handout", "no handout found")
        return
    plan = handouts[0].read_text(encoding="utf8")

    minutes = 0
    for h1, m1, h2, m2 in re.findall(r"\|\s*(\d):(\d\d)[-–](\d):(\d\d)\s*\|", plan):
        minutes += (int(h2) * 60 + int(m2)) - (int(h1) * 60 + int(m1))
    if minutes != 180:
        report.fail("lesson plan",
                    f"segments sum to {minutes} minutes, not 180")

    cited = []
    for start, end in re.findall(r"Slides? (\d+)(?:[-–](\d+))?", plan):
        cited.append(int(start))
        if end:
            cited.append(int(end))
    if cited and max(cited) > slide_count:
        report.fail("lesson plan",
                    f"cites slide {max(cited)} but the deck has {slide_count}")

    check_notes_timings(plan, notes_by_slide, report)


#: A row of the plan: its two clock times and the slides it covers.
PLAN_ROW = re.compile(
    r"\|\s*(\d):(\d\d)[-–](\d):(\d\d)\s*\|[^|]*\|\s*Slides? (\d+)(?:[-–](\d+))?")
#: "22 minutes", "20 minutes." - what a notebook slide's notes tell the lecturer.
STATED = re.compile(r"\b(\d{1,3})\s+minutes\b")


def check_notes_timings(plan: str, notes_by_slide: dict, report: Report) -> None:
    """A slide that owns a whole segment must not tell the lecturer another time.

    Only single-slide segments are checked - the notebooks and the break. Inside
    a segment covering fifteen slides, "two or three minutes" is about something
    else entirely, and reading it as a claim about the segment would make this
    check cry wolf until it was ignored.

    The failure it exists for: two minutes moved out of Notebook 01 to pay for a
    new slide, the plan updated, and the speaker notes left saying 22. Nothing
    could see it. The plan summed to 180 either way.
    """
    for h1, m1, h2, m2, first, last in PLAN_ROW.findall(plan):
        if last and int(last) != int(first):
            continue
        budget = (int(h2) * 60 + int(m2)) - (int(h1) * 60 + int(m1))
        stated = STATED.search(notes_by_slide.get(int(first), ""))
        if stated and int(stated.group(1)) != budget:
            report.fail("lesson plan",
                        f"slide {first}'s notes say {stated.group(1)} minutes, "
                        f"the plan gives that segment {budget}")


# --------------------------------------------------------------- acronyms

def normalise(text: str) -> str:
    """Lower-case, hyphens and underscores to spaces, whitespace collapsed."""
    cleaned = text.replace("*", " ").replace("-", " ").replace("_", " ")
    return " ".join(cleaned.split()).lower()


def strip_maths(text: str) -> str:
    """Remove formulas and code spans before looking for acronyms.

    An acronym inside a formula is notation, not prose: TPR appearing in a
    display equation does not introduce the term to anybody. The sentence
    around it does, and that sentence is what this check is about.
    """
    text = DISPLAY_MATH.sub(" ", text)
    text = INLINE_MATH.sub(" ", text)
    return re.sub(r"`[^`]*`", " ", text)


def check_acronyms(lesson: Path, report: Report) -> None:
    for path in artefacts(lesson):
        body = strip_maths(markdown_of(path))
        for acronym, expansions in ACRONYMS.items():
            match = re.search(rf"\b{acronym}\b", body)
            if not match:
                continue
            line_end = body.find("\n", match.end())
            prefix = normalise(body[:line_end if line_end != -1 else len(body)])
            if not any(normalise(e) in prefix for e in expansions):
                report.fail("acronyms",
                            f"{path.name} uses {acronym} before expanding it")


# -------------------------------------------------------------------- quiz

#: Question stems that are legitimately imperative rather than interrogative.
#: Half of lesson 5's are ("Describe k-fold cross-validation..."), so the check
#: has to accept them or it becomes noise nobody reads.
IMPERATIVES = ("define", "state", "describe", "list", "explain", "compute",
               "summarise", "summarize", "draw", "give", "write", "identify",
               "distinguish", "derive", "compare", "name", "read", "take",
               "prove", "argue", "diagnose", "prescribe", "justify", "sketch",
               "show", "contrast", "resolve", "criticise", "rank", "report",
               "propose", "outline", "discuss", "interpret", "decide", "choose")

def asks_something(stem: str) -> bool:
    """Does this stem actually pose a question or give an instruction?

    Deliberately generous. Most stems set a scene first and ask afterwards, so
    the ask is looked for anywhere rather than at either end. What this rejects
    is a stem that describes a situation and never gets round to the question.
    """
    if "?" in stem:
        return True
    lowered = stem.lower()
    return any(re.search(rf"\b{verb}\b", lowered) for verb in IMPERATIVES)


def question_stem(source: str) -> str:
    """The first paragraph of a question cell, with markup removed.

    Not "everything up to the next **": a stem containing bold would be cut off
    at the first one, which is how this check first reported four false
    positives in lesson 1.
    """
    paragraph = source.strip().split("\n\n")[0]
    paragraph = re.sub(r"^\*\*\d+\.\s*", "", paragraph.strip())

    # Order matters. The whole stem is bold, so its closing asterisks sit after
    # the difficulty marker; strip the markup first or the anchored pattern
    # below never reaches the end of the string.
    paragraph = re.sub(r"<[^>]+>", " ", paragraph)
    paragraph = paragraph.replace("*", " ").replace("`", " ")
    paragraph = re.sub(r"\(\s*(reasoning|recall|calculation)\s*\)\s*$", "",
                       " ".join(paragraph.split()), flags=re.I)
    return " ".join(paragraph.split())


def check_quiz(lesson: Path, report: Report) -> None:
    quizzes = sorted((lesson / "Quizzes").glob("*.ipynb"))
    if not quizzes:
        report.fail("quiz", "no quiz found")
        return

    for path in quizzes:
        name = path.name
        try:
            nb = json.loads(path.read_text(encoding="utf8"))
        except json.JSONDecodeError:
            continue        # check_notebooks already reported it

        numbers: list[int] = []
        sections = 0

        for cell in nb["cells"]:
            source = "".join(cell.get("source", []))
            stripped = source.lstrip()

            if stripped.startswith("## "):
                sections += 1

            starts = re.findall(r"^\*\*(\d+)\.", source, re.M)
            if not starts:
                continue
            if len(starts) > 1:
                report.fail("quiz",
                            f"{name}: one cell holds questions "
                            f"{', '.join(starts)} — each gets its own cell")
            numbers.append(int(starts[0]))

            # The collapsible answer, in the house format.
            if source.count("<details>") != source.count("</details>"):
                report.fail("quiz",
                            f"{name} question {starts[0]}: <details> is not closed")
            if source.count("<details>") != 1:
                report.fail("quiz",
                            f"{name} question {starts[0]}: expected exactly one "
                            f"answer block, found {source.count('<details>')}")
            if "<summary>" not in source or "color='darkgreen'" not in source:
                report.fail("quiz",
                            f"{name} question {starts[0]}: the answer summary is "
                            "not in the course format")

            # An answer that says nothing is worse than no answer, because it
            # looks answered. Length rather than bullet count: a question that
            # asks for the rule "in one sentence" deserves one bullet.
            points = re.findall(r"<li>(.*?)</li>", source, re.S)
            answered = sum(len(re.sub(r"<[^>]+>", "", p).strip()) for p in points)
            if not points:
                report.fail("quiz",
                            f"{name} question {starts[0]}: the answer has no points")
            elif answered < 80:
                report.fail("quiz",
                            f"{name} question {starts[0]}: the answer is a stub "
                            f"({answered} characters)")

            stem = question_stem(source)
            if stem and not asks_something(stem):
                report.fail("quiz",
                            f"{name} question {starts[0]}: the stem states a "
                            f"situation but never asks anything — {stem[:52]}")

        if numbers != list(range(1, len(numbers) + 1)):
            report.fail("quiz",
                        f"{name}: questions are numbered {numbers[:8]}... — "
                        "they must run consecutively from 1")
        if sections == 0:
            report.fail("quiz", f"{name}: no '## ' section headers")
        if len(numbers) < 15:
            report.fail("quiz",
                        f"{name}: only {len(numbers)} questions; a lesson's quiz "
                        "carries at least 15")


# ------------------------------------------------------- worked arithmetic

def check_worked_examples(lesson: Path, report: Report) -> None:
    """Run the lesson's own recomputation of every hand-worked number."""
    checker = lesson / "Docs" / "worked_examples.py"
    handouts = sorted((lesson / "Docs").glob("*.md"))
    handout = handouts[0].read_text(encoding="utf8") if handouts else ""

    # Does the handout work any arithmetic by hand at all?
    works_numbers = bool(
        re.search(r"worked example|by hand|A worked step", handout, re.I))

    if not checker.exists():
        if works_numbers:
            report.fail("worked examples",
                        "the handout works numbers by hand but there is no "
                        "Docs/worked_examples.py recomputing them")
        return

    result = run_checker(checker, report)
    if result.returncode != 0:
        tail = (result.stderr or result.stdout).strip().splitlines()[-4:]
        report.fail("worked examples", " | ".join(tail))
    else:
        for line in result.stdout.strip().splitlines():
            report.note(line)


#: Where docker-compose bind-mounts this repository inside the image.
WORKDIR = "/home/jovyan/work"


def run_checker(checker: Path, report: Report):
    """Recompute the lesson's numbers on the stack that produced them.

    The container first, and the host only as a fallback. Every figure a
    handout quotes came out of a notebook run in the image, and the image
    pins numpy, pandas and scikit-learn; the host venv does not track those
    pins and in practice runs well ahead of them. Most of these checks are
    arithmetic on the dataset and do not care, but the ones that fit a model
    do: the same logistic pipeline scores 0.751401 in the image and 0.751528
    on the host, so a fourth decimal verified here is not the fourth decimal
    a student will see.

    A missing or stopped container is not a failure - it is a note, and the
    checks still run - because verifying on the wrong stack is worth far more
    than not verifying at all.
    """
    # Ask whether the container is there before asking it to do anything, so
    # that "docker is not running" and "your arithmetic is wrong" cannot arrive
    # as the same failure. Matching on the daemon's wording would be one more
    # string to keep in step with docker.
    try:
        alive = subprocess.run(["docker", "exec", CONTAINER, "true"],
                               capture_output=True, text=True).returncode == 0
    except FileNotFoundError:
        alive = False  # no docker installed here
    if alive:
        # The path has to be the one the container sees, not the one we hold:
        # the repository is bind-mounted, the same file under another name.
        relative = checker.resolve().relative_to(ROOT).as_posix()
        return subprocess.run(
            ["docker", "exec", "-w", WORKDIR, CONTAINER, "python", relative],
            capture_output=True, text=True)

    report.note("container not running; the arithmetic ran on the host stack, "
                "which does not match the image's pinned versions")
    return subprocess.run([sys.executable, str(checker)],
                          capture_output=True, text=True, cwd=str(ROOT))


# ------------------------------------------------------- cross-references

def check_cross_references(lesson: Path, report: Report) -> None:
    notebooks = {p.name[:2] for p in (lesson / "Notebooks").glob("*.ipynb")}
    for path in artefacts(lesson):
        body = markdown_of(path)
        for number in re.findall(r"[Nn]otebook (\d)\b", body):
            if f"0{number}" not in notebooks:
                report.fail("cross-references",
                            f"{path.name} mentions notebook {number}, "
                            "which does not exist")


# ------------------------------------------------------- quoted numbers

#: Words too common in this course to identify a claim on their own. A quotation
#: has to reproduce several of a claim's words, and at least one of them has to
#: be one the source lesson uses far more than this one does, or "accuracy" and
#: "model" would tie every lesson to every other one.
COMMON_WORDS = set("""
a an and are as at be been before but by can does for from give gives had has
have how in into is it its just left like made make more most much no not of
off on one only or our out over same than that the their them then there these
they this to two under up was were what when where which who why with without
you your number numbers accuracy model models data set sets score scores test
tests train training lesson lessons section sections notebook figure table
slide slides course student students first last next
""".split())

#: How much more often a word must occur in the lesson being quoted than in the
#: one doing the quoting before it counts as that lesson's word. "Imputed" is
#: lesson 2's; "layers" is not lesson 9's, whatever its carry-home sentence
#: says, because lesson 10 uses it on every other slide.
CLAIM_WORD_RATIO = 3.0

#: What counts as "beside the number": the clause around it, not the paragraph.
#: Widening this to a sentence starts matching a synthesis list - which quotes
#: half a dozen lessons in one breath - against every claim it contains.
QUOTE_BEFORE, QUOTE_AFTER = 60, 110

CARRY_HOME_HEADING = "One number per lesson worth remembering"
TABLE_ROW = re.compile(r"^\|\s*(\d{1,2})\s*\|\s*(.+?)\s*\|\s*$", re.M)
BOLD = re.compile(r"\*\*(.+?)\*\*")
QUANTITY = r"(\d[\d,]*(?:\.\d+)?)\s*(%?)"
NUMBER = re.compile(r"(?<![\w.,])" + QUANTITY)
#: The two-number shapes the carry-home numbers come in: "94 of 128",
#: "0.885 to 1.000", "0.941 against -0.046".
PAIRED = re.compile(r"(?<![\w.,])" + QUANTITY + r"\s+(of|out of|to|against)\s+"
                    + QUANTITY, re.I)
CLAIM_WORD = re.compile(r"[a-z][a-z-]{2,}")


def numbers_in(text: str) -> set[str]:
    return {m.group(0).replace(",", "") for m in
            re.finditer(r"\d[\d,]*(?:\.\d+)?", text)}


def says(value: str, pool) -> bool:
    """Is `value` one of `pool`, allowing for the rounding a quote does?

    CLAUDE.md remembers lesson 10 by "0.857 to 0.462" where the notebooks print
    0.8567 and 0.4617. Quoting a number to three places rather than four is not
    a lesson quoting a different number.
    """
    if value in pool:
        return True
    if "." not in value:
        return False
    places = len(value.split(".")[1])
    for other in pool:
        try:
            if f"{float(other):.{places}f}" == value:
                return True
        except ValueError:
            pass
    return False


def shape_of(value: str, percent: str) -> tuple:
    """What a number looks like: 0.857 and 0.462 match, 3,500 and 0.5 do not."""
    return bool(percent), "." in value, len(value.split(".")[0])


@lru_cache(maxsize=None)
def lesson_prose(lesson: Path) -> str:
    """Every word a student reads in this lesson, markup flattened away."""
    joined = "\n".join(markdown_of(p) for p in artefacts(lesson))
    return " ".join(re.sub(r"[|*`#>]", " ", joined).split())


@lru_cache(maxsize=None)
def lesson_numbers(lesson: Path) -> frozenset:
    """Every number this lesson prints anywhere a student can read it."""
    return frozenset(numbers_in(lesson_prose(lesson)))


@lru_cache(maxsize=None)
def word_rates(lesson: Path) -> dict:
    """How often each word of the course's claims occurs in this lesson.

    Per ten thousand words, so the ten lessons compare despite their different
    lengths.
    """
    prose = lesson_prose(lesson).lower()
    total = max(len(prose.split()), 1) / 10_000
    return {w: len(re.findall(rf"\b{re.escape(w)}\b", prose)) / total
            for claim in carry_home_claims() for w in claim["words"]}


@lru_cache(maxsize=None)
def carry_home_claims() -> tuple:
    """The course's register of quotable numbers: CLAUDE.md's own table.

    Each row names the one number a lesson is meant to be remembered by, which
    makes it exactly the number the *other* lessons repeat - the synthesis
    lists in lessons 8, 9 and 10 are built from them. So the table is what a
    quotation elsewhere has to agree with.
    """
    body = (ROOT / "CLAUDE.md").read_text(encoding="utf8")
    if CARRY_HOME_HEADING not in body:
        return ()
    table = body.split(CARRY_HOME_HEADING, 1)[1].split("\n### ", 1)[0]
    claims = []
    for row in TABLE_ROW.finditer(table):
        lesson, cell = row.group(1).zfill(2), row.group(2)
        if not numbers_in(cell):
            continue                       # the header, and any prose row
        text = " ".join(re.sub(r"[*`]", "", cell).split())
        words = {w for w in CLAIM_WORD.findall(text.lower())
                 if w not in COMMON_WORDS}
        pair = PAIRED.search(text)
        primary = sorted(numbers_in(" ".join(BOLD.findall(cell))))
        claims.append({
            "lesson": lesson,
            "text": text,
            "primary": primary,
            "numbers": numbers_in(cell),
            "words": words,
            # ("94", "of", "128") where the claim is a pair, else None
            "pair": (pair.group(1).replace(",", ""), pair.group(3).lower(),
                     pair.group(4).replace(",", "")) if pair else None,
            "percent": bool(pair is None and len(primary) == 1
                            and re.search(re.escape(primary[0]) + r"\s*%", text)),
        })
    return tuple(claims)


def check_quoted_numbers(lesson: Path, report: Report) -> None:
    """A number quoted from another lesson must still match its source.

    Lessons 8, 9 and 10 each close on a list of the numbers the course has
    produced so far, and all three said "98 of 128 imputed rows" after the
    review of lesson 2 established the figure is 94 - the neighbour search that
    produced 98 was over three standardised columns where the imputer used four
    raw ones. Nothing caught it. The three lessons agreed with each other;
    `worked_examples.py` recomputes only what its *own* handout prints; and
    lesson 2 was correct everywhere it spoke for itself.

    So this compares a quotation against its source rather than against its
    neighbours, in the two directions that can go stale:

    * **the source moved and the table did not** - a lesson no longer contains
      the number CLAUDE.md says it is remembered by;
    * **the source moved and a quotation did not** - another lesson reproduces
      the claim's wording beside a number that is not the claim's, and that
      appears nowhere in the lesson it is taken from.

    Three conditions together are what keep it quiet enough to be worth having.
    A quotation has to reproduce the claim's *wording*, not merely sit near a
    number; the number has to have the claim's shape, so a count is never read
    as a mis-stated accuracy; and it must appear nowhere in the source lesson,
    which is what separates a stale quote from an honest new measurement.
    Lesson 10 runs a network of lesson 9's design on its own wafers and prints
    numbers lesson 9 never saw: those are not quotations of anything.
    """
    claims = carry_home_claims()
    if not claims:
        return
    here = lesson.name[:2]
    lessons = {p.name[:2]: p for p in sorted(ROOT.glob("Lessons/[0-9][0-9]_*"))}
    if here not in lessons:
        return
    source_numbers = {c["lesson"]: lesson_numbers(lessons[c["lesson"]])
                      for c in claims if c["lesson"] in lessons}
    every_claim_number = set().union(*(c["numbers"] for c in claims))

    # Half one: the number this lesson is remembered by is still in it.
    for claim in claims:
        if claim["lesson"] != here:
            continue
        gone = [n for n in claim["primary"]
                if not says(n, source_numbers.get(here, frozenset()))]
        if gone:
            report.fail(
                "quoted numbers",
                f"CLAUDE.md remembers this lesson by \"{claim['text']}\", and "
                f"{', '.join(gone)} appears in no artefact of it - the table is "
                "quoting a number the lesson has moved on from")

    # Half two: what this lesson quotes from the others.
    mine = word_rates(lesson)

    def matches(claim: dict, window: str, needed: int) -> bool:
        """Does this window reproduce the claim, in the source lesson's words?"""
        theirs = word_rates(lessons[claim["lesson"]])
        hits = {w for w in claim["words"]
                if re.search(rf"\b{re.escape(w)}\b", window)}
        return len(hits) >= needed and any(
            theirs[w] and theirs[w] >= CLAIM_WORD_RATIO * mine[w] for w in hits)

    for path in artefacts(lesson):
        flat = " ".join(re.sub(r"[|*`#>]", " ", markdown_of(path)).split())
        low = flat.lower()
        for claim in claims:
            source = claim["lesson"]
            if source == here or source not in source_numbers:
                continue

            def stale(value: str, at: tuple[int, int], needed: int) -> None:
                if value in every_claim_number:
                    return
                if says(value, source_numbers[source]):
                    return              # a number of that lesson's own
                window = low[max(0, at[0] - QUOTE_BEFORE):at[1] + QUOTE_AFTER]
                if not matches(claim, window, needed):
                    return
                report.fail(
                    "quoted numbers",
                    f"{path.name} quotes {value} where lesson {int(source)} "
                    f"gives \"{claim['text']}\", and {value} is in no artefact "
                    "of that lesson - check the source before reprinting it")

            if claim["pair"]:
                # One half of the pair still matches, the other does not: the
                # shape of a quotation left behind by a correction.
                first, joiner, second = claim["pair"]
                for found in PAIRED.finditer(flat):
                    if found.group(3).lower() != joiner:
                        continue
                    left = found.group(1).replace(",", "")
                    right = found.group(4).replace(",", "")
                    if (left == first) == (right == second):
                        continue
                    stale(right if left == first else left, found.span(), 2)
            elif len(claim["primary"]) == 1:
                target = claim["primary"][0]
                wanted = shape_of(target, "%" if claim["percent"] else "")
                for found in NUMBER.finditer(flat):
                    value = found.group(1).replace(",", "")
                    if value == target:
                        continue
                    if shape_of(value, found.group(2)) != wanted:
                        continue
                    stale(value, found.span(), 2)


# -------------------------------------------------------------------- main

#: A slide is 7.5in tall; rendered at 100 dpi that is 563 rows. Text below this
#: row is inside the footer band: it collides with the university crest, and a
#: few rows further down it is simply cut off by the edge of the slide. The
#: value is measured, not derived - across the ten decks every slide either
#: stops by row 522 or runs to 548 and beyond, with nothing in between.
DECK_TEXT_FLOOR = 530

#: The crest sits in this column band. Its own dark lettering would otherwise
#: register as text on every slide.
CREST_COLUMNS = (845, 960)

#: A tighter box around the crest, and the amount of its own dark ink. The
#: crest is identical on every slide: 544 of the course's 555 slides measure
#: exactly CREST_INK here, so anything meaningfully above it is a line of text
#: sitting on top of the logo. A bullet can end below the floor check above and
#: still collide here, which is how lesson 8's t-SNE slide got through.
CREST_BOX = (830, 924)
CREST_ROWS = (487, 551)
CREST_INK = 236
CREST_TOLERANCE = 40


def check_deck_overflow(lesson: Path, report: Report) -> None:
    """Look at the rendered deck and find text that runs off the slide.

    The estimate in postprocess_pptx counts characters, which cannot tell a
    bullet that wraps to two lines from one that just fits: it scored thirteen
    clipped slides exactly as it scored their fixed versions. Rasterising the
    built PDF measures what a student will actually see, which is the whole
    point of the rule in CLAUDE.md about looking at what you built.

    Silently skipped where the tools are missing - it must not become a reason
    the verifier cannot run.
    """
    decks = list((lesson / "Slides").glob("*_slides.pdf"))
    if not decks or shutil.which("pdftoppm") is None:
        return
    try:
        import numpy as np
        from PIL import Image
    except ImportError:
        return

    for deck in decks:
        with tempfile.TemporaryDirectory() as tmp:
            rendered = subprocess.run(
                ["pdftoppm", "-png", "-r", "100", str(deck), f"{tmp}/slide"],
                capture_output=True,
            )
            if rendered.returncode != 0:
                return
            for page in sorted(Path(tmp).glob("slide-*.png")):
                image = np.array(Image.open(page).convert("RGB")).astype(int)
                ink = (image.max(axis=2) < 120) & (
                    image.max(axis=2) - image.min(axis=2) < 40
                )
                crest = ink[CREST_ROWS[0]:CREST_ROWS[1],
                            CREST_BOX[0]:CREST_BOX[1]].sum()
                ink[:, CREST_COLUMNS[0]:CREST_COLUMNS[1]] = False
                rows = np.where(ink.any(axis=1))[0]
                number = int(page.stem.split("-")[-1])
                if len(rows) and rows.max() > DECK_TEXT_FLOOR:
                    report.fail(
                        "deck",
                        f"{deck.name}: slide {number} runs past the bottom of"
                        " the slide - shorten it or split it",
                    )
                elif crest > CREST_INK + CREST_TOLERANCE:
                    report.fail(
                        "deck",
                        f"{deck.name}: slide {number} has text running into"
                        " the university crest - shorten it or split it",
                    )


def verify(lesson: Path, run: bool) -> Report:
    report = Report()
    check_notebooks(lesson, report, run)
    check_pins(lesson, report)
    check_figures(lesson, report)
    check_slides(lesson, report)
    check_deck_overflow(lesson, report)
    check_quiz(lesson, report)
    check_acronyms(lesson, report)
    check_worked_examples(lesson, report)
    check_cross_references(lesson, report)
    check_quoted_numbers(lesson, report)
    return report


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("lesson", nargs="?", help="lesson number, e.g. 05")
    parser.add_argument("--run", action="store_true",
                        help="also execute every notebook in the container")
    args = parser.parse_args()

    lessons = sorted(ROOT.glob("Lessons/[0-9][0-9]_*"))
    if args.lesson:
        lessons = [p for p in lessons if p.name.startswith(args.lesson)]
        if not lessons:
            print(f"no lesson matching {args.lesson}")
            return 1
        # A lesson with no handout is dropped from the sweep below, which is
        # right when verifying everything: one nobody has started yet is not a
        # defect. Asked for by name it is the opposite - lesson 10, stopped
        # halfway through phase A, reported "0 verificate, nessun problema" and
        # exited 0, which reads as a pass on work that does not exist.
        unbuilt = [p for p in lessons if not any(p.glob("Docs/*.md"))]
        if unbuilt:
            for path in unbuilt:
                print(f"{path.name}  —  NON COSTRUITA: nessuna dispensa in "
                      f"Docs/, quindi non c'è niente da verificare.")
            print("\n    .  python .claude/skills/new-lesson/scripts/"
                  f"lesson_state.py {args.lesson}  dice in che fase è.")
            return 1
    else:
        lessons = [p for p in lessons if any(p.glob("Docs/*.md"))]

    failed = 0
    for lesson in lessons:
        report = verify(lesson, args.run)
        status = "OK" if report.ok else f"{len(report.problems)} PROBLEMS"
        print(f"\n{lesson.name}  —  {status}")
        for problem in report.problems:
            print(f"    x  {problem}")
        for note in report.notes:
            print(f"    .  {note}")
        if not report.ok:
            failed += 1

    print()
    if failed:
        print(f"{failed} lesson(s) with problems")
        return 1
    print(f"{len(lessons)} lesson(s) verified, no problems")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
