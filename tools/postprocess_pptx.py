"""Repair a pandoc-generated .pptx so it renders identically everywhere.

    python tools/postprocess_pptx.py deck.pptx [deck2.pptx ...]

Called automatically by tools/build.py after each conversion. Two fixes, in
order, because the second cannot run until the first has made the file valid.

1. **Undeclared namespace prefixes.**
   When pandoc puts an equation inside a table cell it emits `<a14:m>` without
   declaring the `a14` prefix on the slide root, producing XML that is simply
   not well formed. PowerPoint silently tolerates it; stricter readers - lxml,
   and therefore python-pptx - refuse the file outright, and some viewers drop
   the slide. Equations outside tables are emitted correctly, which is why the
   problem only shows up on some slides.

2. **Inherited geometry.**
   pandoc emits shapes with no <a:off>/<a:ext>, leaving position and size to be
   inherited from layout and master. That is valid OOXML and PowerPoint honours
   it, but several viewers fall back to (0, 0) and stack the title on top of the
   body. Course material gets opened in PowerPoint, Google Slides, LibreOffice,
   the OneDrive web preview and on phones, so we resolve the inheritance here
   and write the result into each slide.
"""

import re
import shutil
import sys
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_A14 = "http://schemas.microsoft.com/office/drawing/2010/main"
NS_M = "http://schemas.openxmlformats.org/officeDocument/2006/math"
NS_MC = "http://schemas.openxmlformats.org/markup-compatibility/2006"

# Prefixes pandoc may use without declaring them, and their correct URIs.
NAMESPACES = {
    "a14": "http://schemas.microsoft.com/office/drawing/2010/main",
    "m": "http://schemas.openxmlformats.org/officeDocument/2006/math",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
    "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
}

SLIDE_XML = re.compile(r"ppt/slides/slide\d+\.xml$")


def _root_span(xml: str) -> tuple[int, int] | None:
    """Character span of the root element's opening tag.

    Skips the XML declaration and any processing instruction or comment, so we
    do not append attributes inside `<?xml ... ?>`.
    """
    for match in re.finditer(r"<(?![?!])[^>]*>", xml):
        return match.start(), match.end() - 1
    return None


def _declare_missing(xml: str) -> tuple[str, list[str]]:
    """Add xmlns declarations for any prefix used but not declared."""
    span = _root_span(xml)
    if span is None:
        return xml, []
    start, end = span
    root = xml[start:end]

    added = []
    for prefix, uri in NAMESPACES.items():
        if not re.search(rf"<{prefix}:|\s{prefix}:", xml):
            continue
        if f"xmlns:{prefix}=" in root:
            continue
        root += f' xmlns:{prefix}="{uri}"'
        added.append(prefix)

    if not added:
        return xml, []
    return xml[:start] + root + xml[end:], added


def repair_namespaces(path: Path) -> dict[str, list[str]]:
    """Rewrite the package, declaring namespaces pandoc left dangling."""
    repaired: dict[str, list[str]] = {}
    temp = path.with_suffix(".repair.pptx")

    with zipfile.ZipFile(path) as src, zipfile.ZipFile(
        temp, "w", zipfile.ZIP_DEFLATED
    ) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if SLIDE_XML.match(item.filename):
                fixed, added = _declare_missing(data.decode("utf8"))
                if added:
                    repaired[item.filename.split("/")[-1]] = added
                    data = fixed.encode("utf8")
            dst.writestr(item, data)

    if repaired:
        temp.replace(path)
    else:
        temp.unlink()
    return repaired


def _omath_text(omath) -> str:
    """Plain-text rendering of an OMML equation, for the fallback branch."""
    parts = [node.text or "" for node in omath.iter(f"{{{NS_M}}}t")]
    return "".join(parts).strip() or "[equation]"


def _wrap_equations(tree) -> int:
    """Wrap bare <a14:m> in <mc:AlternateContent>. Returns equations wrapped.

    pandoc drops <a14:m> straight into <a:p>, but CT_TextParagraph only admits
    a:pPr, a:r, a:br, a:fld and a:endParaRPr. Consumers are entitled to reject
    the paragraph, and they do: LibreOffice renders the slide body empty and
    PowerPoint offers to "repair" the file, discarding content as it goes.

    The conformant form puts the extension behind a compatibility gate, with a
    plain-text fallback for readers that cannot render OMML.
    """
    wrapped = 0
    for math in list(tree.iter(f"{{{NS_A14}}}m")):
        parent = math.getparent()
        if parent is None or parent.tag == f"{{{NS_MC}}}Choice":
            continue  # already wrapped

        alt = etree.Element(f"{{{NS_MC}}}AlternateContent", nsmap={"mc": NS_MC})
        choice = etree.SubElement(
            alt, f"{{{NS_MC}}}Choice", nsmap={"a14": NS_A14}, Requires="a14"
        )
        fallback = etree.SubElement(alt, f"{{{NS_MC}}}Fallback")

        run = etree.SubElement(fallback, f"{{{NS_A}}}r")
        etree.SubElement(run, f"{{{NS_A}}}rPr")
        text = etree.SubElement(run, f"{{{NS_A}}}t")
        text.text = _omath_text(math)

        parent.replace(math, alt)
        choice.append(math)
        wrapped += 1
    return wrapped


def gate_equations(path: Path) -> int:
    """Rewrite every slide so its equations are schema-conformant."""
    total = 0
    temp = path.with_suffix(".gate.pptx")

    with zipfile.ZipFile(path) as src, zipfile.ZipFile(
        temp, "w", zipfile.ZIP_DEFLATED
    ) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if SLIDE_XML.match(item.filename):
                tree = etree.fromstring(data)
                count = _wrap_equations(tree)
                if count:
                    total += count
                    data = etree.tostring(
                        tree, xml_declaration=True, encoding="UTF-8", standalone=True
                    )
            dst.writestr(item, data)

    if total:
        temp.replace(path)
    else:
        temp.unlink()
    return total


def _layout_counterpart(shape, layout):
    """The layout placeholder a slide placeholder inherits geometry from."""
    if not shape.is_placeholder:
        return None
    idx = shape.placeholder_format.idx
    for candidate in layout.placeholders:
        if candidate.placeholder_format.idx == idx:
            return candidate
    return None


def pin_geometry(path: Path) -> int:
    """Give every shape explicit position and size. Returns shapes changed."""
    prs = Presentation(str(path))
    pinned = 0

    for slide in prs.slides:
        layout = slide.slide_layout
        for shape in slide.shapes:
            spPr = shape._element.find(qn("p:spPr"))
            if spPr is None:
                continue
            xfrm = spPr.find(qn("a:xfrm"))
            if xfrm is not None and xfrm.find(qn("a:off")) is not None:
                continue  # already explicit

            counterpart = _layout_counterpart(shape, layout)
            resolved = {}
            for attr in ("left", "top", "width", "height"):
                value = getattr(shape, attr)
                if value is None and counterpart is not None:
                    value = getattr(counterpart, attr)
                resolved[attr] = value

            if any(v is None for v in resolved.values()):
                # Nothing to inherit from; leave it rather than invent a place.
                continue

            for attr, value in resolved.items():
                setattr(shape, attr, value)
            pinned += 1

    if pinned:
        prs.save(str(path))
    return pinned


def renumber_oversized_ids(path: Path) -> int:
    """Bring shape ids back under 2**31. Returns how many were rewritten.

    PowerPoint refuses ids at or above 2**31 in ordinary shapes: it opens the
    file with "PowerPoint found a problem with content", repairs it, and drops
    whatever it could not read. Tooling reaches that range easily, because slide
    masters number their layout list from 2**31 upward and a naive "max id plus
    one" walks straight into it.

    Cheap to check, expensive to miss, so we verify every part on every build.
    """
    rewritten = 0
    temp = path.with_suffix(".ids.pptx")

    with zipfile.ZipFile(path) as src, zipfile.ZipFile(
        temp, "w", zipfile.ZIP_DEFLATED
    ) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            is_shape_part = item.filename.startswith(
                ("ppt/slides/slide", "ppt/slideLayouts/", "ppt/slideMasters/")
            ) and item.filename.endswith(".xml")

            if is_shape_part:
                tree = etree.fromstring(data)
                nodes = [
                    node
                    for node in tree.iter()
                    if etree.QName(node).localname == "cNvPr"
                    and (node.get("id") or "").isdigit()
                ]
                used = {int(n.get("id")) for n in nodes if int(n.get("id")) < 2**31}
                changed = False
                for node in nodes:
                    if int(node.get("id")) < 2**31:
                        continue
                    candidate = 1
                    while candidate in used:
                        candidate += 1
                    used.add(candidate)
                    node.set("id", str(candidate))
                    rewritten += 1
                    changed = True
                if changed:
                    data = etree.tostring(
                        tree, xml_declaration=True, encoding="UTF-8", standalone=True
                    )
            dst.writestr(item, data)

    if rewritten:
        temp.replace(path)
    else:
        temp.unlink()
    return rewritten


#: Characters that fit on one full-width line at the body's default size, and
#: the height one line occupies as a fraction of the slide. Both measured from
#: the generated decks rather than derived - pptx does no layout of its own, so
#: there is nothing exact to compute here.
CHARS_PER_LINE = 52
LINE_HEIGHT = 0.085


def _text_height(shape, slide_h: int) -> int:
    """Estimate how tall a text box needs to be for its contents.

    A fixed height is wrong as soon as a slide carries more than a line or two:
    the text overflows its shape, the figure beneath is positioned against the
    shape rather than the visible text, and the two overlap.
    """
    lines = 0
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        lines += max(1, -(-len(text) // CHARS_PER_LINE))  # ceil division
    return int(slide_h * LINE_HEIGHT * max(lines, 1))


def relayout_figure_slides(path: Path) -> int:
    """Give slides that carry a picture a full-width title and centred figure.

    Any slide holding both text and an image makes pandoc reach for the
    "Content with Caption" layout, which is a two-column design: a narrow
    caption column on the left, content on the right. For a lecture slide
    showing an equation or a plot that reads badly - the title shrinks to a
    third of the width and drops to 15pt while everything else stays large.

    We reflow those slides to the shape the material actually wants: title
    across the top at full size, supporting text beneath it, figure centred in
    the space that remains, aspect ratio preserved.
    """
    prs = Presentation(str(path))
    slide_w, slide_h = prs.slide_width, prs.slide_height
    margin = int(slide_w * 0.05)
    usable = slide_w - 2 * margin
    changed = 0

    for slide in prs.slides:
        pictures = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
        if not pictures:
            continue
        title = slide.shapes.title
        if title is None:
            continue

        if title.width < usable * 0.9:
            # A caption layout pandoc left in its narrow two-column form.
            title.left, title.top = margin, int(slide_h * 0.04)
            title.width, title.height = usable, int(slide_h * 0.17)

        top = title.top + title.height + int(slide_h * 0.02)
        # Compare the underlying XML elements: python-pptx builds a fresh proxy
        # object on every access, so `is not title` would never match and the
        # title itself would get shuffled down with the body shapes.
        bodies = [
            s
            for s in slide.shapes
            if s.has_text_frame
            and s._element is not title._element
            and s.text_frame.text.strip()
        ]
        for body in bodies:
            body.left, body.top = margin, top
            body.width = usable
            body.height = _text_height(body, slide_h)
            top += body.height + int(slide_h * 0.02)

        available_h = slide_h - top - int(slide_h * 0.12)  # keep clear of logo
        for picture in pictures:
            ratio = picture.width / picture.height
            height = min(available_h, int(usable / ratio))
            width = int(height * ratio)
            picture.width, picture.height = width, height
            picture.left = int((slide_w - width) / 2)
            picture.top = top + int((available_h - height) / 2)
        changed += 1

    if changed:
        prs.save(str(path))
    return changed


def process(path: Path, verbose: bool = True) -> bool:
    backup = path.with_suffix(".orig.pptx")
    shutil.copy2(path, backup)
    try:
        repaired = repair_namespaces(path)
        equations = gate_equations(path)
        pinned = pin_geometry(path)
        reflowed = relayout_figure_slides(path)
        renumbered = renumber_oversized_ids(path)
    except Exception as exc:  # keep the original rather than a broken file
        shutil.copy2(backup, path)
        print(f"  {path.name}: postprocess FAILED ({exc}) - left as generated")
        return False
    finally:
        backup.unlink(missing_ok=True)

    if verbose:
        detail = []
        if repaired:
            detail.append(f"namespaces fixed on {len(repaired)} slide(s)")
        if equations:
            detail.append(f"{equations} equation(s) gated")
        if pinned:
            detail.append(f"{pinned} shape(s) pinned")
        if reflowed:
            detail.append(f"{reflowed} figure slide(s) reflowed")
        if renumbered:
            detail.append(f"{renumbered} oversized id(s) renumbered")
        print(f"  {path.name}: {'; '.join(detail) if detail else 'nothing to fix'}")
    return True


def main(argv: list[str]) -> int:
    if not argv:
        print(__doc__, file=sys.stderr)
        return 2
    ok = True
    for name in argv:
        path = Path(name)
        if not path.exists():
            print(f"missing: {path}", file=sys.stderr)
            return 1
        ok &= process(path)
    return 0 if ok else 1


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
