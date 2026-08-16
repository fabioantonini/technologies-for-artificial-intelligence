"""Build Course/template.pptx, the pandoc reference document for all slide decks.

Run this only when the visual identity changes. The generated template is
committed, so day-to-day slide building does not depend on this script.

    python tools/make_template.py

The template starts from pandoc's own default reference document, so every
placeholder pandoc expects stays intact. We only change the theme colours,
the fonts and add the university logo.
"""

import re
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

import postprocess_pptx  # noqa: E402
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.oxml.shapes.picture import CT_Picture
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent
LOGO = ROOT / "Course" / "univaq_logo.png"
OUT = ROOT / "Course" / "template.pptx"

# Palette. Gold is sampled from the university crest; the rest is a neutral
# academic scheme chosen so that body text keeps a high contrast ratio.
GOLD = "F0BE50"  # crest gold, accent only - never used for body text
INK = "1A1A1A"  # titles and body text
BLUE = "1F4E79"  # primary accent
TEAL = "2E7D8A"
SLATE = "4A5568"
RUST = "9C4221"
HEADING_FONT = "Segoe UI Semibold"
BODY_FONT = "Segoe UI"


def fetch_pandoc_reference(dest: Path) -> None:
    """Ask pandoc for its default pptx reference document."""
    with dest.open("wb") as fh:
        subprocess.run(
            ["pandoc", "--print-default-data-file", "reference.pptx"],
            stdout=fh,
            check=True,
        )


def restyle_theme(pptx_path: Path) -> None:
    """Rewrite theme colours and fonts inside the packaged theme XML."""
    tmp = pptx_path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(pptx_path) as src, zipfile.ZipFile(
        tmp, "w", zipfile.ZIP_DEFLATED
    ) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if re.match(r"ppt/theme/theme\d+\.xml$", item.filename):
                data = _patch_theme(data.decode("utf8")).encode("utf8")
            dst.writestr(item, data)
    tmp.replace(pptx_path)


def _patch_theme(xml: str) -> str:
    accents = [BLUE, TEAL, GOLD, SLATE, RUST, INK]
    for i, colour in enumerate(accents, start=1):
        xml = re.sub(
            rf'(<a:accent{i}>\s*<a:srgbClr val=")[0-9A-Fa-f]{{6}}',
            rf"\g<1>{colour}",
            xml,
        )
    xml = re.sub(
        r'(<a:dk1>\s*<a:sysClr val="windowText" lastClr=")[0-9A-Fa-f]{6}',
        rf"\g<1>{INK}",
        xml,
    )
    xml = re.sub(
        r'(<a:majorFont>\s*<a:latin typeface=")[^"]*', rf"\g<1>{HEADING_FONT}", xml
    )
    xml = re.sub(
        r'(<a:minorFont>\s*<a:latin typeface=")[^"]*', rf"\g<1>{BODY_FONT}", xml
    )
    return xml


def _safe_shape_id(shapes) -> int:
    """Smallest unused shape id in this shape tree.

    We cannot use python-pptx's `_next_shape_id` on a slide master. It looks up
    ids with an XPath beginning `//`, which lxml resolves from the document
    root rather than from the shape tree, so it also picks up the entries in
    <p:sldLayoutIdLst> - and those start at 2**31 by convention. The id it
    hands back is then above 2**31 too, which PowerPoint rejects: the file
    opens with "PowerPoint found a problem with content" and gets repaired,
    losing whatever it could not read. Layouts have no such list, which is why
    only the master was affected.
    """
    used = {
        int(value)
        for value in shapes._spTree.xpath(".//@id")
        if value.isdigit() and int(value) < 2**31
    }
    candidate = 1
    while candidate in used:
        candidate += 1
    return candidate


def _place_picture(base_slide, image_path, left, top, width, height):
    """Add a picture to a master or layout.

    python-pptx only exposes add_picture on SlideShapes, so we drive the same
    underlying machinery by hand: register the image part, then append a
    <p:pic> to the shape tree.
    """
    image_part, r_id = base_slide.part.get_or_add_image_part(str(image_path))
    shapes = base_slide.shapes
    shape_id = _safe_shape_id(shapes)
    pic = CT_Picture.new_pic(
        shape_id,
        f"Logo {shape_id}",
        image_part.desc,
        r_id,
        Emu(int(left)),
        Emu(int(top)),
        Emu(int(width)),
        Emu(int(height)),
    )
    shapes._spTree.append(pic)


def normalise_caption_layouts(pptx_path: Path) -> None:
    """Make the caption layouts look like the ordinary ones.

    pandoc reaches for "Content with Caption" whenever a slide holds both text
    and a picture, and "Picture with Caption" when it holds only a picture.
    Both ship as narrow two-column designs whose title is pinned to 15pt in a
    third of the width - fine for a photo essay, wrong for a lecture slide
    showing an equation or a plot, and jarring next to every other slide.

    We widen the title to full width and drop the hard-coded font sizes so the
    text inherits from the master, exactly like "Title and Content" does.
    tools/postprocess_pptx.py then centres the figure in what is left.
    """
    prs = Presentation(str(pptx_path))
    slide_w, slide_h = prs.slide_width, prs.slide_height
    margin = int(slide_w * 0.05)
    usable = slide_w - 2 * margin

    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name not in ("Content with Caption", "Picture with Caption"):
                continue
            for shape in layout.placeholders:
                kind = shape.placeholder_format.type
                if kind == PP_PLACEHOLDER.TITLE:
                    shape.left, shape.top = margin, int(slide_h * 0.04)
                    shape.width, shape.height = usable, int(slide_h * 0.17)
                elif kind in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.PICTURE):
                    shape.left, shape.top = margin, int(slide_h * 0.24)
                    shape.width = usable
                    shape.height = int(slide_h * 0.16)

                # Drop hard-coded sizes so the master's styles apply.
                for prop in shape._element.iter(qn("a:defRPr"), qn("a:rPr")):
                    prop.attrib.pop("sz", None)
    prs.save(str(pptx_path))


def add_logos(pptx_path: Path) -> None:
    """Place the crest on the title layout and a small mark on the master."""
    prs = Presentation(str(pptx_path))
    slide_w, slide_h = prs.slide_width, prs.slide_height
    ratio = 342 / 784  # logo aspect ratio

    for master in prs.slide_masters:
        # Small footer mark, bottom right, inherited by every ordinary slide.
        w = slide_w * 0.115
        h = w * ratio
        _place_picture(
            master,
            LOGO,
            slide_w - w - slide_w * 0.028,
            slide_h - h - slide_h * 0.030,
            w,
            h,
        )
        for layout in master.slide_layouts:
            if layout.name != "Title Slide":
                continue
            # The opening slide carries the crest at full size, so suppress the
            # inherited footer mark rather than showing the logo twice.
            layout._element.set("showMasterSp", "0")
            w = slide_w * 0.26
            h = w * ratio
            _place_picture(layout, LOGO, (slide_w - w) / 2, slide_h * 0.055, w, h)
    prs.save(str(pptx_path))


def main() -> int:
    if not LOGO.exists():
        print(f"missing logo: {LOGO}", file=sys.stderr)
        return 1
    OUT.parent.mkdir(parents=True, exist_ok=True)
    fetch_pandoc_reference(OUT)
    restyle_theme(OUT)
    normalise_caption_layouts(OUT)
    add_logos(OUT)

    # pandoc's reference document ships four sample slides that inherit their
    # geometry. They never reach a generated deck - only the masters, layouts
    # and theme are read from here - but the template is a file people open, so
    # it should not greet them with overlapping placeholders or a repair prompt.
    postprocess_pptx.process(OUT, verbose=False)

    size = OUT.stat().st_size / 1024
    print(f"written {OUT.relative_to(ROOT)} ({size:.0f} KB)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
